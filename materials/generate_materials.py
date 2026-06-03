"""Generate Word guide and PPT slides from Markdown source files."""

from pathlib import Path


def _add_code_block(document, lines):
    paragraph = document.add_paragraph()
    run = paragraph.add_run('\n'.join(lines))
    run.font.name = 'Consolas'


def _add_markdown_table(document, rows):
    parsed = [
        [cell.strip() for cell in row.strip().strip('|').split('|')]
        for row in rows
        if not set(row.replace('|', '').replace(' ', '').strip()) <= {'-', ':'}
    ]
    if not parsed:
        return
    table = document.add_table(rows=len(parsed), cols=len(parsed[0]))
    table.style = 'Table Grid'
    for row_index, row in enumerate(parsed):
        for column_index, value in enumerate(row):
            table.cell(row_index, column_index).text = value


def generate_word():
    from docx import Document

    source = Path('materials/teacher_lab_guide.md')
    document = Document()
    in_code = False
    code_lines = []
    table_lines = []

    def flush_table():
        nonlocal table_lines
        if table_lines:
            _add_markdown_table(document, table_lines)
            table_lines = []

    for line in source.read_text(encoding='utf-8').splitlines():
        if line.startswith('```'):
            flush_table()
            if in_code:
                _add_code_block(document, code_lines)
                code_lines = []
                in_code = False
            else:
                in_code = True
            continue
        if in_code:
            code_lines.append(line)
            continue
        if line.strip().startswith('|') and line.strip().endswith('|'):
            table_lines.append(line)
            continue
        flush_table()
        if line.startswith('# '):
            document.add_heading(line[2:], level=0)
        elif line.startswith('## '):
            document.add_heading(line[3:], level=1)
        elif line.startswith('### '):
            document.add_heading(line[4:], level=2)
        elif line.startswith('- '):
            document.add_paragraph(line[2:], style='List Bullet')
        elif len(line) > 3 and line[0].isdigit() and '. ' in line[:5]:
            document.add_paragraph(line.split('. ', 1)[1], style='List Number')
        elif line.strip():
            document.add_paragraph(line)
    flush_table()
    if code_lines:
        _add_code_block(document, code_lines)

    target = Path('materials/teacher_lab_guide.docx')
    try:
        document.save(target)
        print(f'Generated {target}')
    except PermissionError:
        fallback = Path('materials/teacher_lab_guide_student_v2.docx')
        document.save(fallback)
        print(f'{target} 正在被 Word/WPS 占用，已生成备用文件：{fallback}')


def _parse_slides():
    slides = []
    current = None
    for line in Path('materials/lecture_slides_outline.md').read_text(encoding='utf-8').splitlines():
        if line.startswith('## Slide'):
            if current:
                slides.append(current)
            title = line.split('：', 1)[-1] if '：' in line else line.replace('## ', '')
            current = {'title': title, 'bullets': []}
        elif current and line.startswith('- '):
            current['bullets'].append(line[2:])
    if current:
        slides.append(current)
    return slides


def generate_ppt():
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    slides = _parse_slides()
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    navy = RGBColor(15, 38, 84)
    blue = RGBColor(37, 99, 235)
    green = RGBColor(16, 133, 92)
    amber = RGBColor(217, 119, 6)
    light = RGBColor(241, 245, 249)
    dark = RGBColor(30, 41, 59)
    white = RGBColor(255, 255, 255)

    def add_box(slide, left, top, width, height, fill):
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = fill

    def add_text(slide, left, top, width, height, text, size, color, bold=False, align=None):
        shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        shape.text_frame.margin_left = Inches(0.05)
        shape.text_frame.margin_right = Inches(0.05)
        shape.text_frame.margin_top = Inches(0.02)
        shape.text_frame.margin_bottom = Inches(0.02)
        shape.text = text
        paragraph = shape.text_frame.paragraphs[0]
        if align is not None:
            paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = 'Microsoft YaHei'
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color

    total = len(slides)
    for index, item in enumerate(slides):
        slide = presentation.slides.add_slide(blank)
        if index == 0:
            add_box(slide, 0, 0, 13.333, 7.5, navy)
            add_text(slide, 0.8, 1.1, 11.7, 0.5, item['bullets'][0], 22, RGBColor(125, 211, 252), True, PP_ALIGN.CENTER)
            add_text(slide, 0.8, 2.1, 11.7, 0.8, item['bullets'][1], 38, white, True, PP_ALIGN.CENTER)
            add_text(slide, 0.8, 3.15, 11.7, 0.5, item['bullets'][2], 22, light, True, PP_ALIGN.CENTER)
            add_text(slide, 0.8, 5.8, 11.7, 0.35, item['bullets'][3], 16, light, False, PP_ALIGN.CENTER)
            continue

        add_box(slide, 0, 0, 13.333, 0.72, navy)
        add_text(slide, 0.5, 0.15, 10.5, 0.35, item['title'], 22, white, True)
        add_text(slide, 11.5, 0.18, 1.2, 0.3, f'{index + 1:02d}/{total:02d}', 12, white, False, PP_ALIGN.RIGHT)
        add_box(slide, 0.55, 1.08, 5.7, 5.7, light)
        accent_colors = [blue, green, amber, RGBColor(14, 165, 233)]
        for bullet_index, bullet in enumerate(item['bullets'][:6]):
            y = 1.35 + bullet_index * 0.82
            add_box(slide, 0.85, y, 0.12, 0.46, accent_colors[bullet_index % len(accent_colors)])
            add_text(slide, 1.1, y + 0.03, 4.8, 0.36, bullet, 15, dark, bullet_index == 0)

        add_box(slide, 6.8, 1.08, 5.95, 5.7, RGBColor(248, 250, 252))
        keywords = ['分集', 'SC', 'MRC', 'DSSS', 'PN', '处理增益']
        for tag_index, tag in enumerate(keywords):
            x = 7.2 + (tag_index % 3) * 1.75
            y = 1.6 + (tag_index // 3) * 0.75
            add_box(slide, x, y, 1.25, 0.38, accent_colors[tag_index % len(accent_colors)])
            add_text(slide, x + 0.05, y + 0.08, 1.15, 0.18, tag, 10, white, True, PP_ALIGN.CENTER)
        if 'MRC' in item['title']:
            visual = 'y = sum(conj(h) * r) / sum(abs(h)^2)'
        elif 'DSSS' in item['title'] or '扩频' in item['title']:
            visual = 'bit -> symbol x PN -> correlation -> bit'
        elif 'SC' in item['title'] or '选择' in item['title']:
            visual = 'k = argmax |h|^2, y = r[k] / h[k]'
        else:
            visual = '实验03：分集可靠性 + 扩频抗干扰'
        add_text(slide, 7.15, 3.45, 5.1, 0.8, visual, 20, navy, True, PP_ALIGN.CENTER)
        add_text(slide, 7.15, 5.4, 5.1, 0.5, 'PR标题：实验03-姓名-学号', 16, amber, True, PP_ALIGN.CENTER)

    target = Path('materials/lecture_slides.pptx')
    presentation.save(target)
    print(f'Generated {target}')


def main():
    generate_word()
    generate_ppt()
    for path in [Path('materials/teacher_lab_guide.docx'), Path('materials/lecture_slides.pptx')]:
        if path.exists():
            print(f'{path}: {path.stat().st_size} bytes')
    temporary_files = list(Path('materials').glob('~$*'))
    if temporary_files:
        raise RuntimeError(f'发现 Word/WPS 临时文件: {temporary_files}')


if __name__ == '__main__':
    main()
